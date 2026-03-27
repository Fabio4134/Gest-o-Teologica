import os
import json
import uvicorn
import requests
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional
import fitz  # PyMuPDF
from dotenv import load_dotenv

# Try to import pptx, but don't fail if not present
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

load_dotenv()

app = FastAPI(title="Question Bank AI Specialist")

# Configurar CORS para permitir o frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Em produção, restringir ao domínio do frontend
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

class QuestionOption(BaseModel):
    text: str
    is_correct: bool

class GeneratedQuestion(BaseModel):
    text: str
    options: List[str]
    correctOptionIndex: int
    explanation: Optional[str] = None

def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_pptx(file_path: str) -> str:
    if not HAS_PPTX:
        raise HTTPException(status_code=500, detail="Suporte para PPTX não disponível no servidor.")
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

@app.post("/generate-questions")
async def generate_questions(
    file: UploadFile = File(...),
    subject_id: str = Form(...),
    quantity: int = Form(5),
    difficulty: str = Form("Médio"),
    types: str = Form("[\"multipla_escolha\"]")
):
    temp_path = None
    try:
        # Salvar arquivo temporariamente
        temp_path = f"temp_{file.filename}"
        with open(temp_path, "wb") as buffer:
            buffer.write(await file.read())

        # Extrair texto
        content_text = ""
        if file.filename.lower().endswith(".pdf"):
            content_text = extract_text_from_pdf(temp_path)
        elif file.filename.lower().endswith(".pptx"):
            content_text = extract_text_from_pptx(temp_path)
        else:
            raise HTTPException(status_code=400, detail="Formato de arquivo não suportado. Use PDF ou PPTX.")

        if len(content_text) < 50:
            raise HTTPException(status_code=400, detail="O arquivo parece conter pouco texto ou não pôde ser lido.")

        # Construir Prompt baseado no docs/prova-online-expert.md
        methodology = ""
        try:
            # Look for docs in parent folder
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            doc_path = os.path.join(base_dir, "docs", "prova-online-expert.md")
            if os.path.exists(doc_path):
                with open(doc_path, "r", encoding="utf-8") as f:
                    methodology = f.read()
        except Exception as e:
            print(f"Erro ao ler metodologia: {e}")

        question_types = json.loads(types)
        
        prompt = f"""
        {methodology}
        
        MISSÃO:
        Gere {quantity} questões de nível {difficulty} baseadas no material fornecido abaixo.
        
        TIPOS DE QUESTÃO SOLICITADOS: {', '.join(question_types)}
        
        MATERIAL DE REFERÊNCIA:
        ---
        {content_text[:6000]}
        ---
        
        Retorne APENAS um JSON válido seguindo esta estrutura:
        {{
          "questions": [
            {{
              "text": "Enunciado",
              "options": ["Opção A", "Opção B", "Opção C", "Opção D"],
              "correctOptionIndex": 0,
              "explanation": "Explicação"
            }}
          ]
        }}
        """

        # Call OpenAI via requests to avoid jiter/pydantic-core issues
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {OPENAI_API_KEY}"
        }
        
        payload = {
            "model": "gpt-4o",
            "messages": [
                {"role": "system", "content": "Você é um assistente especializado em criação de questões teológicas para o sistema IBAD."},
                {"role": "user", "content": prompt}
            ],
            "response_format": { "type": "json_object" }
        }

        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
        
        if response.status_code != 200:
            raise Exception(f"Erro na API OpenAI: {response.text}")

        response_data = response.json()
        result = json.loads(response_data["choices"][0]["message"]["content"])
        
        # Injetar subject_id nas questões
        for q in result.get("questions", []):
            q["subject_id"] = subject_id

        return result

    except Exception as e:
        print(f"Erro generalizado: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except:
                pass

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
